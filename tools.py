import subprocess
import sys
import os
import json
import imaplib
import smtplib
import email
from email.header import decode_header
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.base import MIMEBase
from email import encoders
import requests
import tempfile

# --- Credential store ---
CREDENTIALS_FILE = os.path.join(os.path.dirname(__file__), "credentials.json")

def _load_creds() -> dict:
    try:
        with open(CREDENTIALS_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}

def _save_creds(data: dict):
    with open(CREDENTIALS_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)
from bs4 import BeautifulSoup
from duckduckgo_search import DDGS
from docx import Document
from fpdf import FPDF
import openpyxl

OBSIDIAN_VAULT = "C:/Users/dani1/OneDrive/Documentos/Obsidian Vault"
MEMORY_FILE    = f"{OBSIDIAN_VAULT}/Adam/MEMORY.md"

ISSE_BASE_URL = "https://isse-certificados.fly.dev"

def _gmail_creds():
    c = _load_creds().get("gmail", {})
    return c.get("address", ""), c.get("app_password", "")

TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "run_command",
            "description": "Run a terminal command on the user's Windows PC. Use PowerShell syntax. Good for opening apps, checking system status, managing files.",
            "parameters": {
                "type": "object",
                "properties": {
                    "command": {"type": "string", "description": "PowerShell command to run"},
                },
                "required": ["command"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": "Read contents of a file on the PC. Handles plain text, PDF (.pdf), Word (.docx), PowerPoint (.pptx), and Excel (.xlsx/.xls). Use this on assignment reference materials downloaded from Teams.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "max_chars": {"type": "integer", "description": "Optional cap on returned chars (default 60000)."},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_file",
            "description": "Write or create a file on the PC.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "content": {"type": "string"},
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_directory",
            "description": "List files and folders in a directory.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "web_search",
            "description": "Search the internet using DuckDuckGo. Returns top results with titles, URLs, snippets.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Search query"},
                    "max_results": {"type": "integer", "description": "Number of results (default 5)"},
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "fetch_webpage",
            "description": "Fetch and read the text content of any webpage URL.",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "Full URL to fetch"},
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_emails",
            "description": "Read recent emails from Daniel's Gmail inbox.",
            "parameters": {
                "type": "object",
                "properties": {
                    "count": {"type": "integer", "description": "Number of recent emails to fetch (default 5)"},
                    "folder": {"type": "string", "description": "Folder to read: INBOX, Sent, etc. Default: INBOX"},
                },
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "save_to_memory",
            "description": "Save something to Adam's long-term memory. Use when Daniel says 'remember this', 'save this', or shares important info to keep.",
            "parameters": {
                "type": "object",
                "properties": {
                    "content": {"type": "string", "description": "The fact, info or note to remember. Be concise and clear."},
                },
                "required": ["content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_memory",
            "description": "Read Adam's full long-term memory. Use when Daniel asks what you remember or to check stored info.",
            "parameters": {
                "type": "object",
                "properties": {},
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "obsidian_read_note",
            "description": "Read a note from Daniel's Obsidian vault by name.",
            "parameters": {
                "type": "object",
                "properties": {
                    "note_name": {"type": "string", "description": "Note filename without .md extension, e.g. 'Meeting Notes'"},
                },
                "required": ["note_name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "obsidian_write_note",
            "description": "Create or overwrite a note in Daniel's Obsidian vault.",
            "parameters": {
                "type": "object",
                "properties": {
                    "note_name": {"type": "string", "description": "Note filename without .md extension"},
                    "content": {"type": "string", "description": "Markdown content to write"},
                },
                "required": ["note_name", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "obsidian_append_note",
            "description": "Append text to an existing Obsidian note without overwriting it.",
            "parameters": {
                "type": "object",
                "properties": {
                    "note_name": {"type": "string", "description": "Note filename without .md extension"},
                    "content": {"type": "string", "description": "Text to append"},
                },
                "required": ["note_name", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "obsidian_list_notes",
            "description": "List all notes in Daniel's Obsidian vault or a subfolder.",
            "parameters": {
                "type": "object",
                "properties": {
                    "subfolder": {"type": "string", "description": "Optional subfolder path within vault, e.g. 'Projects/Petroleum'"},
                },
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "obsidian_search_notes",
            "description": "Search for a keyword across all notes in Daniel's Obsidian vault.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Keyword to search for"},
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "obsidian_create_folder",
            "description": "Create a folder inside Daniel's Obsidian vault.",
            "parameters": {
                "type": "object",
                "properties": {
                    "folder_path": {"type": "string", "description": "Folder path relative to vault root, e.g. 'Projects/Petroleum'"},
                },
                "required": ["folder_path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "obsidian_move_note",
            "description": "Move a note from one location to another inside the Obsidian vault.",
            "parameters": {
                "type": "object",
                "properties": {
                    "source": {"type": "string", "description": "Current note path relative to vault, without .md"},
                    "destination": {"type": "string", "description": "New note path relative to vault, without .md"},
                },
                "required": ["source", "destination"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "obsidian_delete_note",
            "description": "Delete a note from Daniel's Obsidian vault. Ask confirmation before using.",
            "parameters": {
                "type": "object",
                "properties": {
                    "note_path": {"type": "string", "description": "Note path relative to vault, without .md"},
                },
                "required": ["note_path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_word_document",
            "description": "Create a Word (.docx) document with given content and save it to a path.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Full path including filename, e.g. C:/Users/dani1/Desktop/report.docx"},
                    "title": {"type": "string", "description": "Document title"},
                    "content": {"type": "string", "description": "Document body text. Use \\n for new lines."},
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_pdf",
            "description": "Create a PDF document with given content and save it to a path.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Full path including filename, e.g. C:/Users/dani1/Desktop/report.pdf"},
                    "title": {"type": "string", "description": "PDF title"},
                    "content": {"type": "string", "description": "PDF body text"},
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_python_file",
            "description": "Create a Python (.py) script file with the given code.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Full path including filename, e.g. C:/Users/dani1/Desktop/solution.py"},
                    "code": {"type": "string", "description": "Python code to write"},
                },
                "required": ["path", "code"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_jupyter_notebook",
            "description": "Create a Jupyter notebook (.ipynb) file with markdown and code cells.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Full path including filename, e.g. C:/Users/dani1/Desktop/analysis.ipynb"},
                    "cells": {
                        "type": "array",
                        "description": "List of cells. Each cell: {type: 'markdown'|'code', content: 'string'}",
                        "items": {
                            "type": "object",
                            "properties": {
                                "type": {"type": "string"},
                                "content": {"type": "string"}
                            }
                        }
                    },
                },
                "required": ["path", "cells"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_powerpoint",
            "description": "Create a PowerPoint (.pptx) presentation with slides.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Full path including filename, e.g. C:/Users/dani1/Desktop/presentation.pptx"},
                    "title": {"type": "string", "description": "Presentation title"},
                    "slides": {
                        "type": "array",
                        "description": "List of slides. Each slide: {title: 'string', content: 'string'}",
                        "items": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string"},
                                "content": {"type": "string"}
                            }
                        }
                    },
                },
                "required": ["path", "slides"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_text_file",
            "description": "Create any plain text file (.txt, .md, .csv, .json, .html, .css, .js, .sql, .r, .java, .cpp, .c, etc.)",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Full path with correct extension, e.g. C:/Users/dani1/Desktop/report.md"},
                    "content": {"type": "string", "description": "File content"},
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_latex_document",
            "description": "Write a full LaTeX source file and compile it to PDF. Use for academic papers, scientific reports, equations, formal typeset documents. Pass complete LaTeX source (\\documentclass … \\end{document}). Returns path to the compiled PDF.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Output path WITHOUT extension, e.g. C:/Users/dani1/Desktop/report — .tex and .pdf will be created here"},
                    "latex_source": {"type": "string", "description": "Complete LaTeX source code from \\documentclass to \\end{document}"},
                },
                "required": ["path", "latex_source"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_formal_word_document",
            "description": "Create a professionally formatted Word (.docx) document with cover page, table of contents, styled headings, page numbers, and header/footer. Use for formal reports, essays, business documents, proposals.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Full path including filename, e.g. C:/Users/dani1/Desktop/report.docx"},
                    "title": {"type": "string", "description": "Document title"},
                    "author": {"type": "string", "description": "Author name"},
                    "date": {"type": "string", "description": "Date string, e.g. 'April 25, 2026'"},
                    "abstract": {"type": "string", "description": "Abstract or executive summary (optional)"},
                    "sections": {
                        "type": "array",
                        "description": "List of sections. Each section: {heading, content, level} where level=1 main, 2 sub, 3 subsub",
                        "items": {
                            "type": "object",
                            "properties": {
                                "heading": {"type": "string"},
                                "content": {"type": "string"},
                                "level": {"type": "integer"},
                            },
                        },
                    },
                },
                "required": ["path", "title", "sections"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "create_spreadsheet",
            "description": "Create an Excel (.xlsx) spreadsheet with rows of data.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Full path including filename, e.g. C:/Users/dani1/Desktop/data.xlsx"},
                    "headers": {"type": "array", "items": {"type": "string"}, "description": "Column headers"},
                    "rows": {"type": "string", "description": "JSON string of rows, e.g. '[[\"Alice\",100],[\"Bob\",200]]'"},
                },
                "required": ["path", "headers", "rows"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "search_emails",
            "description": "Search emails in Gmail by keyword.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Search term (sender, subject, keyword)"},
                    "count": {"type": "integer", "description": "Max results (default 5)"},
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "save_credential",
            "description": "Save a credential for a service (e.g. teams client_id, any API key). Stored securely in credentials.json.",
            "parameters": {
                "type": "object",
                "properties": {
                    "service": {"type": "string", "description": "Service name, e.g. 'teams', 'gmail', 'notion', 'spotify'"},
                    "key": {"type": "string", "description": "Credential key, e.g. 'client_id', 'api_key', 'password'"},
                    "value": {"type": "string", "description": "The credential value"},
                },
                "required": ["service", "key", "value"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_credentials",
            "description": "List all stored services and their credential keys (not values) so Daniel knows what is configured.",
            "parameters": {"type": "object", "properties": {}, "required": []},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "desktop_screenshot",
            "description": "Take a screenshot of the current screen or a specific window. Returns a vision image Sonnet can analyze to see what is on screen and decide next steps.",
            "parameters": {
                "type": "object",
                "properties": {
                    "window_title": {"type": "string", "description": "Optional: partial window title to capture only that window, e.g. 'Teams'. Leave empty for full screen."},
                },
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "app_focus_window",
            "description": "Bring a window to the foreground by partial title match.",
            "parameters": {
                "type": "object",
                "properties": {
                    "title": {"type": "string", "description": "Partial window title, e.g. 'Teams', 'Chrome', 'Notepad'"},
                },
                "required": ["title"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "app_click",
            "description": "Click at screen coordinates (x, y). Use desktop_screenshot first to identify coordinates.",
            "parameters": {
                "type": "object",
                "properties": {
                    "x": {"type": "integer", "description": "X coordinate"},
                    "y": {"type": "integer", "description": "Y coordinate"},
                    "button": {"type": "string", "description": "'left' or 'right' (default left)"},
                    "double": {"type": "boolean", "description": "Double click (default false)"},
                },
                "required": ["x", "y"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "app_type",
            "description": "Type text into the currently focused element.",
            "parameters": {
                "type": "object",
                "properties": {
                    "text": {"type": "string", "description": "Text to type"},
                },
                "required": ["text"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "app_hotkey",
            "description": "Press a keyboard shortcut, e.g. 'ctrl+g', 'ctrl+c', 'alt+tab', 'enter', 'escape'.",
            "parameters": {
                "type": "object",
                "properties": {
                    "keys": {"type": "string", "description": "Key combo, e.g. 'ctrl+g' or 'enter' or 'alt+tab'"},
                },
                "required": ["keys"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "app_scroll",
            "description": "Scroll the mouse wheel at screen coordinates. Use to scroll inside Teams panels, chat feeds, or assignment lists.",
            "parameters": {
                "type": "object",
                "properties": {
                    "x": {"type": "integer", "description": "X coordinate to scroll at"},
                    "y": {"type": "integer", "description": "Y coordinate to scroll at"},
                    "direction": {"type": "string", "description": "'down' or 'up'"},
                    "clicks": {"type": "integer", "description": "Number of scroll clicks (default 5)"},
                },
                "required": ["x", "y", "direction"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "app_click_text",
            "description": "Click a UI element based on its visible text within a specific window. More reliable than clicking coordinates.",
            "parameters": {
                "type": "object",
                "properties": {
                    "text": {"type": "string", "description": "The visible text of the element to click (e.g., a button label)."},
                    "window_title": {"type": "string", "description": "The partial title of the window containing the element (e.g., 'Microsoft Teams')."},
                },
                "required": ["text", "window_title"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_bot_code",
            "description": "Read Adam's own source code files (main.py, tools.py, call.py). Use to understand current implementation before modifying.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file": {"type": "string", "description": "File to read: 'main.py', 'tools.py', or 'call.py'"},
                },
                "required": ["file"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "edit_bot_code",
            "description": "Edit Adam's own source code. Replaces old_string with new_string in the specified file. Always read_bot_code first to see current state.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file": {"type": "string", "description": "File to edit: 'main.py', 'tools.py', or 'call.py'"},
                    "old_string": {"type": "string", "description": "Exact string to replace (must match exactly)"},
                    "new_string": {"type": "string", "description": "Replacement string"},
                },
                "required": ["file", "old_string", "new_string"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "restart_bot",
            "description": "Restart Adam's bot process to apply code changes. Use after editing main.py or tools.py.",
            "parameters": {"type": "object", "properties": {}, "required": []},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "rollback_bot_code",
            "description": "Restore a bot source file from its .bak backup. Use if a recent edit caused issues.",
            "parameters": {
                "type": "object",
                "properties": {"file": {"type": "string", "description": "main.py, tools.py, or call.py"}},
                "required": ["file"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_improvement_log",
            "description": "Read recent self-improvement attempts (success/failure history). ALWAYS check this BEFORE editing code — avoid repeating past mistakes.",
            "parameters": {
                "type": "object",
                "properties": {"limit": {"type": "integer", "description": "Number of recent entries (default 20)"}},
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_improvement_log",
            "description": "Log a self-improvement attempt — what you tried, whether it worked, and why. Call this after EVERY improvement attempt (code edit, direct retry, or install) so future rounds can learn from it.",
            "parameters": {
                "type": "object",
                "properties": {
                    "status": {"type": "string", "description": "TRIED_APPROACH / ADDED_TOOL / INSTALLED_PKG / DIRECT_RETRY / GAVE_UP"},
                    "detail": {"type": "string", "description": "What you tried and why it should (or didn't) work"},
                    "snippet": {"type": "string", "description": "Relevant code snippet or tool name (optional)"},
                },
                "required": ["status", "detail"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "save_pending_task",
            "description": "Save the current task AND conversation history so it is re-attempted with full context after bot restarts. Always call this BEFORE calling restart_bot when you are doing a self-improvement code edit.",
            "parameters": {
                "type": "object",
                "properties": {
                    "chat_id": {"type": "integer", "description": "Telegram chat ID to send result to"},
                    "task": {"type": "string", "description": "The original task Daniel requested"},
                    "history": {"type": "array", "description": "Recent conversation history (last 10 messages) for context preservation", "items": {"type": "object"}},
                },
                "required": ["chat_id", "task"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "pip_install",
            "description": "Install a Python package using pip. Use when a task requires a library not yet installed.",
            "parameters": {
                "type": "object",
                "properties": {
                    "package": {"type": "string", "description": "Package name, e.g. 'requests' or 'pandas==2.0.0'"},
                },
                "required": ["package"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "download_file",
            "description": "Download any file from the internet and save it to disk.",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "Direct URL to download"},
                    "save_path": {"type": "string", "description": "Full path to save the file, e.g. C:/Users/dani1/Desktop/file.zip"},
                },
                "required": ["url", "save_path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "teams_launch_app",
            "description": "Launch the Microsoft Teams desktop application.",
            "parameters": {"type": "object", "properties": {}, "required": []},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "send_telegram_file",
            "description": "Send a file or folder to Daniel via Telegram. Use when Daniel asks to 'pass', 'send', or 'share' a file or folder. Folders are zipped automatically.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Full path to the file or folder to send"},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "youtube_play",
            "description": "Search and play a song or video on YouTube. Opens browser, searches, plays first result. Use this for any 'play X on youtube' request.",
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "Song or video to search for, e.g. 'Despacito Luis Fonsi'"},
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "media_control",
            "description": "Control media playback on the PC using Windows media keys. Use for pause, play, next track, previous track, stop — works with any media player (Chrome, Spotify, VLC, etc.).",
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {"type": "string", "description": "One of: play_pause, next, previous, stop"},
                },
                "required": ["action"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "browser_get_interactive_elements",
            "description": "Get all visible inputs and buttons on the current page. Use this to discover what you can click or type into before acting.",
            "parameters": {"type": "object", "properties": {}, "required": []},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "set_volume",
            "description": "Set the system volume on Daniel's PC. Level is 0-100.",
            "parameters": {
                "type": "object",
                "properties": {
                    "level": {"type": "integer", "description": "Volume level 0-100"},
                },
                "required": ["level"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_volume",
            "description": "Get the current system volume level.",
            "parameters": {"type": "object", "properties": {}, "required": []},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "set_brightness",
            "description": "Set screen brightness on Daniel's PC. Level is 0-100. Works on laptops; may not work on external desktop monitors.",
            "parameters": {
                "type": "object",
                "properties": {
                    "level": {"type": "integer", "description": "Brightness level 0-100"},
                },
                "required": ["level"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "get_brightness",
            "description": "Get the current screen brightness level.",
            "parameters": {"type": "object", "properties": {}, "required": []},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "browser_navigate",
            "description": "Open a URL in the browser. Launches browser if not open. Use for YouTube, Google, any website.",
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "Full URL to navigate to"},
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "browser_click",
            "description": "Click an element on the current page. Provide CSS selector or visible text.",
            "parameters": {
                "type": "object",
                "properties": {
                    "selector": {"type": "string", "description": "CSS selector, e.g. 'button#search', or visible text with prefix 'text=Play song'"},
                },
                "required": ["selector"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "browser_type",
            "description": "Type text into an input field. Optionally press Enter after.",
            "parameters": {
                "type": "object",
                "properties": {
                    "selector": {"type": "string", "description": "CSS selector of the input field"},
                    "text": {"type": "string", "description": "Text to type"},
                    "press_enter": {"type": "boolean", "description": "Press Enter after typing (default false)"},
                },
                "required": ["selector", "text"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "browser_screenshot",
            "description": "Take a screenshot of the current browser page. Returns file path. Use to see what's on screen before clicking.",
            "parameters": {
                "type": "object",
                "properties": {},
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "browser_get_content",
            "description": "Get visible text content of the current page. Useful to find element text, read results, check what's loaded.",
            "parameters": {
                "type": "object",
                "properties": {},
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "browser_press_key",
            "description": "Press a keyboard key on the current page, e.g. Enter, Escape, ArrowDown, Tab, Space.",
            "parameters": {
                "type": "object",
                "properties": {
                    "key": {"type": "string", "description": "Key name, e.g. 'Enter', 'Escape', 'Space', 'ArrowDown'"},
                },
                "required": ["key"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "browser_scroll",
            "description": "Scroll the page up or down.",
            "parameters": {
                "type": "object",
                "properties": {
                    "direction": {"type": "string", "description": "'up' or 'down'"},
                    "amount": {"type": "integer", "description": "Pixels to scroll (default 500)"},
                },
                "required": ["direction"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "browser_wait",
            "description": "Wait for an element to appear on the page, or just wait N milliseconds.",
            "parameters": {
                "type": "object",
                "properties": {
                    "selector": {"type": "string", "description": "CSS selector to wait for (optional)"},
                    "ms": {"type": "integer", "description": "Milliseconds to wait if no selector given (default 1000)"},
                },
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "browser_evaluate",
            "description": "Run JavaScript on the current page and return the result. Advanced use.",
            "parameters": {
                "type": "object",
                "properties": {
                    "script": {"type": "string", "description": "JavaScript expression to evaluate"},
                },
                "required": ["script"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "browser_close",
            "description": "Close the browser.",
            "parameters": {
                "type": "object",
                "properties": {},
                "required": [],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "send_email",
            "description": "Send an email from Daniel's Gmail account. Can attach files. Ask for confirmation before sending.",
            "parameters": {
                "type": "object",
                "properties": {
                    "to": {"type": "string", "description": "Recipient email address"},
                    "subject": {"type": "string", "description": "Email subject line"},
                    "body": {"type": "string", "description": "Email body (plain text)"},
                    "cc": {"type": "string", "description": "Optional CC address"},
                    "attachments": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Optional list of full file paths to attach, e.g. ['C:/Users/dani1/Desktop/report.pdf']"
                    },
                },
                "required": ["to", "subject", "body"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "isse_get_ciclos",
            "description": "Get all available academic periods (ciclos) loaded in the ISSE teacher certificate system. Use to know which semesters are available before querying.",
            "parameters": {"type": "object", "properties": {}, "required": []},
        },
    },
    {
        "type": "function",
        "function": {
            "name": "isse_search_professor",
            "description": "Search professors by name in the ISSE system. Returns list of matching professors with IDs. Always call this first before isse_get_certificates.",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string", "description": "Professor name or partial name to search (e.g. 'Garcia', 'Juan')"},
                },
                "required": ["name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "isse_get_certificates",
            "description": "Get teaching certificate records for a professor. Returns courses, hours, departments per academic period.",
            "parameters": {
                "type": "object",
                "properties": {
                    "professor_id": {"type": "string", "description": "Professor name (exact or partial) from isse_search_professor result"},
                    "ciclos": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Optional list of specific ciclos to filter (e.g. ['PERIODO 2024-1']). Leave empty for all periods.",
                    },
                },
                "required": ["professor_id"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "isse_export_certificate",
            "description": "Export a professor's teaching certificate as an xlsx or csv file to Downloads folder. Returns the local file path. After this, use send_telegram_file to send it or send_email to email it.",
            "parameters": {
                "type": "object",
                "properties": {
                    "professor_name": {"type": "string", "description": "Professor name (partial match)"},
                    "ciclos": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Optional ciclo filter list. Leave empty for all periods.",
                    },
                    "fmt": {"type": "string", "enum": ["excel", "csv"], "description": "Export format: 'excel' (default) or 'csv'"},
                },
                "required": ["professor_name"],
            },
        },
    },
]


# --- Safety layer ---

# Hard blocks — never execute under any circumstance
BLOCKED_PATTERNS = [
    "format c:", "format d:", "format e:",           # disk format
    "del /s /q c:\\", "rd /s /q c:\\",               # wipe C drive
    "remove-item -recurse -force c:\\",              # PowerShell wipe
    "system32", "syswow64",                          # system dirs
    "reg delete hklm", "reg delete hkcu",            # registry wipe
    "bcdedit", "bootrec",                            # bootloader
    "cipher /w",                                     # secure wipe
    "netsh firewall set opmode disable",             # disable firewall
    "net user administrator",                        # admin account tamper
    "takeown /f c:\\windows",                        # take ownership of Windows
]

# Risky patterns — no longer blocked, Daniel has full access
RISKY_PATTERNS = []

# Safe directories — full PC access, no restriction
SAFE_DIRS = ["c:/"]  # entire C drive allowed


# Stores any pending action that needs !confirm
pending_action = {"name": None, "args": None}


def is_blocked(command: str) -> str | None:
    cmd_lower = command.lower()
    for pattern in BLOCKED_PATTERNS:
        if pattern in cmd_lower:
            return f"🚫 BLOCKED: Command contains forbidden pattern '{pattern}'. This action is permanently restricted for PC safety."
    return None


def is_risky(command: str) -> bool:
    cmd_lower = command.lower()
    return any(p in cmd_lower for p in RISKY_PATTERNS)


def is_safe_path(path: str) -> bool:
    return any(path.lower().startswith(d) for d in SAFE_DIRS)


# --- Tool implementations ---

def _set_pending(name: str, args: dict) -> str:
    pending_action["name"] = name
    pending_action["args"] = args
    return "PENDING_CONFIRM"


def confirm_pending() -> str:
    name = pending_action.get("name")
    args = pending_action.get("args")
    if not name:
        return "No pending action to confirm."
    pending_action["name"] = None
    pending_action["args"] = None
    return _execute_confirmed(name, args)


def _execute_confirmed(name: str, args: dict) -> str:
    """Execute a previously confirmed action, bypassing safety checks."""
    if name == "run_command":
        return _run_command_direct(args["command"])
    elif name == "obsidian_delete_note":
        return _obsidian_delete_direct(args["note_path"])
    return f"Unknown confirmed action: {name}"


def _run_command_direct(command: str) -> str:
    blocked = is_blocked(command)
    if blocked:
        return blocked
    try:
        result = subprocess.run(
            ["powershell", "-Command", command],
            capture_output=True, text=True, timeout=120
        )
        return (result.stdout or result.stderr or "(no output)")[:2000]
    except subprocess.TimeoutExpired:
        return "Command timed out."
    except Exception as e:
        return f"Error: {e}"


def _obsidian_delete_direct(note_path: str) -> str:
    try:
        path = os.path.join(OBSIDIAN_VAULT, f"{note_path}.md")
        if not os.path.exists(path):
            return f"Note '{note_path}' not found."
        os.remove(path)
        return f"Deleted: {note_path}.md"
    except Exception as e:
        return f"Error: {e}"


def run_command(command: str) -> str:
    blocked = is_blocked(command)
    if blocked:
        return blocked
    return _run_command_direct(command)


def read_file(path: str, max_chars: int = 60000) -> str:
    try:
        ext = os.path.splitext(path)[1].lower()
        if ext == ".pdf":
            import pypdf
            reader = pypdf.PdfReader(path)
            parts = []
            for i, page in enumerate(reader.pages):
                try:
                    parts.append(f"--- Page {i+1} ---\n{page.extract_text() or ''}")
                except Exception as e:
                    parts.append(f"--- Page {i+1} (extract error: {e}) ---")
            return "\n".join(parts)[:max_chars]
        if ext == ".docx":
            import docx
            d = docx.Document(path)
            paras = [p.text for p in d.paragraphs if p.text.strip()]
            for tbl in d.tables:
                for row in tbl.rows:
                    paras.append(" | ".join(c.text for c in row.cells))
            return "\n".join(paras)[:max_chars]
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(path)
            parts = []
            for i, slide in enumerate(prs.slides):
                buf = [f"--- Slide {i+1} ---"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        buf.append(shape.text)
                parts.append("\n".join(buf))
            return "\n".join(parts)[:max_chars]
        if ext in (".xlsx", ".xls"):
            try:
                from openpyxl import load_workbook
                wb = load_workbook(path, data_only=True)
                parts = []
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    parts.append(f"--- Sheet: {sheet} ---")
                    for row in ws.iter_rows(values_only=True):
                        parts.append(" | ".join("" if v is None else str(v) for v in row))
                return "\n".join(parts)[:max_chars]
            except Exception as e:
                return f"Error reading spreadsheet: {e}"
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return f.read(max_chars)
    except Exception as e:
        return f"Error reading file: {type(e).__name__}: {e}"


def write_file(path: str, content: str) -> str:
    try:
        dir_path = os.path.dirname(path)
        if dir_path:
            os.makedirs(dir_path, exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        return f"File written: {path}"
    except Exception as e:
        return f"Error writing file: {e}"


def list_directory(path: str) -> str:
    try:
        entries = os.listdir(path)
        return "\n".join(entries) if entries else "(empty)"
    except Exception as e:
        return f"Error: {e}"


def web_search(query: str, max_results: int = 5) -> str:
    try:
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=max_results))
        if not results:
            return "No results found."
        output = []
        for r in results:
            output.append(f"**{r['title']}**\n{r['href']}\n{r['body']}\n")
        return "\n".join(output)
    except Exception as e:
        return f"Search error: {e}"


def fetch_webpage(url: str) -> str:
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        resp = requests.get(url, headers=headers, timeout=15)
        soup = BeautifulSoup(resp.text, "lxml")
        for tag in soup(["script", "style", "nav", "footer"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        return text[:4000]
    except Exception as e:
        return f"Error fetching page: {e}"


def _decode_str(s):
    if s is None:
        return ""
    parts = decode_header(s)
    decoded = []
    for part, enc in parts:
        if isinstance(part, bytes):
            decoded.append(part.decode(enc or "utf-8", errors="replace"))
        else:
            decoded.append(part)
    return "".join(decoded)


def _connect_gmail():
    addr, pwd = _gmail_creds()
    if not addr or not pwd:
        return None, "Gmail credentials not configured. Ask Adam to save them with save_credential."
    mail = imaplib.IMAP4_SSL("imap.gmail.com")
    mail.login(addr, pwd)
    return mail, None


def get_emails(count: int = 5, folder: str = "INBOX") -> str:
    mail, err = _connect_gmail()
    if err:
        return err
    try:
        mail.select(folder)
        _, data = mail.search(None, "ALL")
        ids = data[0].split()[-count:]
        results = []
        for eid in reversed(ids):
            _, msg_data = mail.fetch(eid, "(RFC822)")
            msg = email.message_from_bytes(msg_data[0][1])
            subject = _decode_str(msg["Subject"])
            sender = _decode_str(msg["From"])
            date = msg["Date"]
            body = ""
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        body = part.get_payload(decode=True).decode("utf-8", errors="replace")[:500]
                        break
            else:
                body = msg.get_payload(decode=True).decode("utf-8", errors="replace")[:500]
            results.append(f"From: {sender}\nDate: {date}\nSubject: {subject}\n{body}\n---")
        mail.logout()
        return "\n".join(results)
    except Exception as e:
        return f"Email error: {e}"


def search_emails(query: str, count: int = 5) -> str:
    mail, err = _connect_gmail()
    if err:
        return err
    try:
        mail.select("INBOX")
        _, data = mail.search(None, f'TEXT "{query}"')
        ids = data[0].split()[-count:]
        if not ids:
            return "No emails found matching that query."
        results = []
        for eid in reversed(ids):
            _, msg_data = mail.fetch(eid, "(RFC822)")
            msg = email.message_from_bytes(msg_data[0][1])
            subject = _decode_str(msg["Subject"])
            sender = _decode_str(msg["From"])
            date = msg["Date"]
            results.append(f"From: {sender}\nDate: {date}\nSubject: {subject}\n---")
        mail.logout()
        return "\n".join(results)
    except Exception as e:
        return f"Search error: {e}"


def save_to_memory(content: str) -> str:
    try:
        from datetime import date
        os.makedirs(os.path.dirname(MEMORY_FILE), exist_ok=True)
        with open(MEMORY_FILE, "a", encoding="utf-8") as f:
            f.write(f"\n- [{date.today()}] {content}\n")
        return f"Saved to long-term memory: {content}"
    except Exception as e:
        return f"Error saving memory: {e}"


def read_memory() -> str:
    try:
        with open(MEMORY_FILE, "r", encoding="utf-8") as f:
            return f.read()
    except FileNotFoundError:
        return "(no long-term memory yet)"
    except Exception as e:
        return f"Error reading memory: {e}"


def obsidian_read_note(note_name: str) -> str:
    try:
        path = os.path.join(OBSIDIAN_VAULT, f"{note_name}.md")
        with open(path, "r", encoding="utf-8") as f:
            return f.read(5000)
    except FileNotFoundError:
        return f"Note '{note_name}' not found in vault."
    except Exception as e:
        return f"Error: {e}"


def obsidian_write_note(note_name: str, content: str) -> str:
    try:
        path = os.path.join(OBSIDIAN_VAULT, f"{note_name}.md")
        os.makedirs(os.path.dirname(path), exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        return f"Note saved: {note_name}.md"
    except Exception as e:
        return f"Error: {e}"


def obsidian_append_note(note_name: str, content: str) -> str:
    try:
        path = os.path.join(OBSIDIAN_VAULT, f"{note_name}.md")
        with open(path, "a", encoding="utf-8") as f:
            f.write(f"\n{content}")
        return f"Appended to: {note_name}.md"
    except Exception as e:
        return f"Error: {e}"


def obsidian_list_notes(subfolder: str = "") -> str:
    try:
        base = os.path.join(OBSIDIAN_VAULT, subfolder) if subfolder else OBSIDIAN_VAULT
        notes = []
        for root, dirs, files in os.walk(base):
            dirs[:] = [d for d in dirs if not d.startswith(".")]
            for f in files:
                if f.endswith(".md"):
                    rel = os.path.relpath(os.path.join(root, f), OBSIDIAN_VAULT)
                    notes.append(rel)
        return "\n".join(notes) if notes else "(no notes found)"
    except Exception as e:
        return f"Error: {e}"


def obsidian_create_folder(folder_path: str) -> str:
    try:
        full_path = os.path.join(OBSIDIAN_VAULT, folder_path)
        os.makedirs(full_path, exist_ok=True)
        # Create a .gitkeep so folder shows in Obsidian
        keep = os.path.join(full_path, ".gitkeep")
        if not os.path.exists(keep):
            open(keep, "w").close()
        return f"Folder created: {folder_path}"
    except Exception as e:
        return f"Error: {e}"


def obsidian_move_note(source: str, destination: str) -> str:
    try:
        src = os.path.join(OBSIDIAN_VAULT, f"{source}.md")
        dst = os.path.join(OBSIDIAN_VAULT, f"{destination}.md")
        os.makedirs(os.path.dirname(dst), exist_ok=True)
        os.rename(src, dst)
        return f"Moved: {source} → {destination}"
    except FileNotFoundError:
        return f"Note '{source}' not found."
    except Exception as e:
        return f"Error: {e}"


def obsidian_delete_note(note_path: str) -> str:
    return _obsidian_delete_direct(note_path)


def obsidian_search_notes(query: str) -> str:
    try:
        results = []
        for root, dirs, files in os.walk(OBSIDIAN_VAULT):
            dirs[:] = [d for d in dirs if not d.startswith(".")]
            for f in files:
                if f.endswith(".md"):
                    path = os.path.join(root, f)
                    with open(path, "r", encoding="utf-8", errors="ignore") as fh:
                        content = fh.read()
                    if query.lower() in content.lower():
                        rel = os.path.relpath(path, OBSIDIAN_VAULT)
                        snippet = content[max(0, content.lower().find(query.lower())-50):][:200]
                        results.append(f"**{rel}**\n...{snippet}...\n")
        return "\n".join(results[:10]) if results else "No notes found matching query."
    except Exception as e:
        return f"Error: {e}"


# --- Credential management ---

def save_credential(service: str, key: str, value: str) -> str:
    data = _load_creds()
    if service not in data:
        data[service] = {}
    data[service][key] = value
    _save_creds(data)
    return f"Saved {service}.{key}"

def list_credentials() -> str:
    data = _load_creds()
    if not data:
        return "No credentials stored yet."
    lines = []
    for service, keys in data.items():
        key_list = ", ".join(k for k in keys if keys[k])  # only non-empty
        empty = ", ".join(k for k in keys if not keys[k])
        line = f"[{service}] configured: {key_list or 'none'}"
        if empty:
            line += f" | missing: {empty}"
        lines.append(line)
    return "\n".join(lines)


# --- Microsoft Teams (Graph API + MSAL) ---

# --- Desktop / App control ---

def _find_window(title_partial: str):
    import win32gui, win32con
    result = []
    def cb(hwnd, _):
        if win32gui.IsWindowVisible(hwnd):
            t = win32gui.GetWindowText(hwnd)
            if title_partial.lower() in t.lower():
                result.append(hwnd)
    win32gui.EnumWindows(cb, None)
    return result[0] if result else None

def desktop_screenshot(window_title: str = "") -> str:
    try:
        import pyautogui
        from PIL import ImageGrab
        import win32gui, win32con

        path = os.path.join(tempfile.gettempdir(), "adam_desktop.png")

        if window_title:
            hwnd = _find_window(window_title)
            if hwnd:
                win32gui.ShowWindow(hwnd, win32con.SW_RESTORE)
                win32gui.SetForegroundWindow(hwnd)
                import time; time.sleep(0.5)
                rect = win32gui.GetWindowRect(hwnd)
                img = ImageGrab.grab(bbox=rect)
                img.save(path)
                return f"VISION_SCREENSHOT:{path}"

        # Full screen fallback
        img = pyautogui.screenshot()
        img.save(path)
        return f"VISION_SCREENSHOT:{path}"
    except Exception as e:
        return f"Screenshot error: {e}"

def app_focus_window(title: str) -> str:
    try:
        import win32gui, win32con
        hwnd = _find_window(title)
        if not hwnd:
            return f"Window '{title}' not found."
        win32gui.ShowWindow(hwnd, win32con.SW_RESTORE)
        win32gui.SetForegroundWindow(hwnd)
        import time; time.sleep(0.4)
        return f"Focused: {win32gui.GetWindowText(hwnd)}"
    except Exception as e:
        return f"Error: {e}"

def app_click(x: int, y: int, button: str = "left", double: bool = False) -> str:
    try:
        import pyautogui
        if double:
            pyautogui.doubleClick(x, y, button=button)
        else:
            pyautogui.click(x, y, button=button)
        return f"Clicked ({x}, {y})"
    except Exception as e:
        return f"Click error: {e}"

def app_type(text: str) -> str:
    try:
        import pyautogui, time
        time.sleep(0.2)
        pyautogui.write(text, interval=0.04)
        return f"Typed: {text}"
    except Exception as e:
        return f"Type error: {e}"

def app_hotkey(keys: str) -> str:
    try:
        import pyautogui, time
        time.sleep(0.2)
        parts = [k.strip() for k in keys.lower().split("+")]
        if len(parts) == 1:
            pyautogui.press(parts[0])
        else:
            pyautogui.hotkey(*parts)
        time.sleep(0.3)
        return f"Pressed: {keys}"
    except Exception as e:
        return f"Hotkey error: {e}"

def app_scroll(x: int, y: int, direction: str, clicks: int = 5) -> str:
    try:
        import pyautogui, time
        pyautogui.moveTo(x, y)
        time.sleep(0.1)
        amount = -clicks if direction == "down" else clicks
        pyautogui.scroll(amount, x=x, y=y)
        time.sleep(0.4)
        return f"Scrolled {direction} {clicks} clicks at ({x}, {y})"
    except Exception as e:
        return f"Scroll error: {e}"

def app_click_text(text: str, window_title: str) -> str:
    """Clicks a UI element by its text using pywinauto."""
    try:
        from pywinauto import Application, findwindows
        import win32gui, win32con
        import time

        # First, focus the window
        hwnd = _find_window(window_title)
        if not hwnd:
            return f"Window '{window_title}' not found."
        
        try:
            win32gui.ShowWindow(hwnd, win32con.SW_RESTORE)
            win32gui.SetForegroundWindow(hwnd)
            time.sleep(0.5)
        except Exception as e:
            return f"Error focusing window: {e}"

        # Connect with pywinauto
        try:
            app = Application(backend="uia").connect(handle=hwnd, timeout=20)
            win = app.window(handle=hwnd)
            win.wait('visible', timeout=10)
        except findwindows.TimeoutError:
            return f"Timed out connecting to window '{window_title}'. Is it responsive?"

        # Find the control by its text
        try:
            # More robustly find the control by iterating through descendants
            ctrls = win.descendants(title=text)
            if not ctrls:
                raise findwindows.ElementNotFoundError
            
            # Click the first visible and enabled control
            for ctrl in ctrls:
                if ctrl.is_visible() and ctrl.is_enabled():
                    ctrl.click_input()
                    return f"Clicked element with text '{text}'."
            
            # If no control was clickable
            return f"Element '{text}' found but it's not visible or enabled."

        except findwindows.ElementNotFoundError:
            return f"Element with text '{text}' not found in window '{window_title}'."
        except Exception as e:
            return f"Unexpected error while searching for element: {type(e).__name__}: {e}"

    except ImportError:
        return "Error: pywinauto is not installed. Please run: pip install pywinauto"
    except Exception as e:
        return f"An unexpected error occurred in app_click_text: {type(e).__name__}: {e}"



# --- Microsoft Teams ---

BOT_DIR = os.path.dirname(os.path.abspath(__file__))
EDITABLE_FILES = {"main.py", "tools.py", "call.py"}
IMPROVEMENT_LOG = os.path.join(BOT_DIR, "improvement_log.jsonl")

def _log_improvement(file: str, status: str, detail: str, snippet: str = ""):
    """Append a structured entry to the improvement log."""
    import datetime, json as _json
    entry = {
        "ts": datetime.datetime.now().isoformat(timespec="seconds"),
        "file": file,
        "status": status,
        "detail": detail[:500],
        "snippet": snippet[:200],
    }
    try:
        with open(IMPROVEMENT_LOG, "a", encoding="utf-8") as f:
            f.write(_json.dumps(entry, ensure_ascii=False) + "\n")
    except Exception:
        pass

def read_improvement_log(limit: int = 20) -> str:
    """Return the last N improvement log entries — what was tried, what worked, what failed."""
    if not os.path.exists(IMPROVEMENT_LOG):
        return "(improvement log empty — no past attempts yet)"
    try:
        with open(IMPROVEMENT_LOG, "r", encoding="utf-8") as f:
            lines = f.readlines()
        recent = lines[-limit:]
        return "\n".join(line.strip() for line in recent)
    except Exception as e:
        return f"Log error: {e}"

def read_bot_code(file: str) -> str:
    if file not in EDITABLE_FILES:
        return f"Only these files are readable: {EDITABLE_FILES}"
    try:
        path = os.path.join(BOT_DIR, file)
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    except Exception as e:
        return f"Error reading {file}: {e}"

def _extract_function_names(source: str) -> list:
    """Return all top-level function names defined in source."""
    import ast
    try:
        tree = ast.parse(source)
        return [n.name for n in ast.walk(tree) if isinstance(n, (ast.FunctionDef, ast.AsyncFunctionDef))]
    except Exception:
        return []

def _validate_python(source: str) -> tuple[bool, str]:
    """Returns (ok, error_msg). Checks syntax + duplicate function defs."""
    import ast
    try:
        ast.parse(source)
    except SyntaxError as e:
        return False, f"SyntaxError on line {e.lineno}: {e.msg}"

    names = _extract_function_names(source)
    dupes = [n for n in set(names) if names.count(n) > 1]
    if dupes:
        return False, f"Duplicate function definitions: {dupes}. Each function must be unique."
    return True, ""

def edit_bot_code(file: str, old_string: str, new_string: str) -> str:
    if file not in EDITABLE_FILES:
        return f"Only these files are editable: {EDITABLE_FILES}"
    try:
        path = os.path.join(BOT_DIR, file)
        with open(path, "r", encoding="utf-8") as f:
            content = f.read()
        if old_string not in content:
            return f"Error: old_string not found in {file}. Read the file first to get exact text."

        new_content = content.replace(old_string, new_string, 1)

        # GUARD 1: Syntax + duplicate check before write
        if file.endswith(".py"):
            ok, err = _validate_python(new_content)
            if not ok:
                # Log to improvement log so Sonnet learns
                _log_improvement(file, "REJECTED", f"edit_bot_code rejected: {err}", new_string[:200])
                return (
                    f"❌ EDIT REJECTED: {err}\n"
                    "Read the file again, fix the issue, and retry. "
                    "If a function with that name already exists, use edit_bot_code to MODIFY it "
                    "instead of adding a duplicate."
                )

        # Backup
        backup = path + ".bak"
        with open(backup, "w", encoding="utf-8") as f:
            f.write(content)

        # Write new
        with open(path, "w", encoding="utf-8") as f:
            f.write(new_content)

        # GUARD 2: post-write import smoke test (only for tools.py)
        if file == "tools.py":
            test = subprocess.run(
                [sys.executable, "-c", "import importlib.util,sys; "
                 "spec=importlib.util.spec_from_file_location('t', r'%s'); "
                 "m=importlib.util.module_from_spec(spec); spec.loader.exec_module(m)" % path],
                capture_output=True, text=True, timeout=15
            )
            if test.returncode != 0:
                # ROLLBACK
                with open(path, "w", encoding="utf-8") as f:
                    f.write(content)
                _log_improvement(file, "ROLLED_BACK", f"import error: {test.stderr[-300:]}", new_string[:200])
                return f"❌ Import test failed — ROLLED BACK. Error:\n{test.stderr[-400:]}\nFix and retry."

        _log_improvement(file, "EDITED_OK", f"replaced {len(old_string)} chars", new_string[:200])
        return f"✅ Edited {file}. Syntax + import OK. Backup at {file}.bak. Call restart_bot to apply."
    except Exception as e:
        return f"Error editing {file}: {e}"

def rollback_bot_code(file: str) -> str:
    """Restore a file from its .bak backup."""
    if file not in EDITABLE_FILES:
        return f"Only these files: {EDITABLE_FILES}"
    path = os.path.join(BOT_DIR, file)
    backup = path + ".bak"
    if not os.path.exists(backup):
        return f"No backup found for {file}"
    try:
        with open(backup, "r", encoding="utf-8") as f:
            content = f.read()
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        _log_improvement(file, "ROLLBACK_MANUAL", "restored from .bak", "")
        return f"✅ Rolled back {file} from {file}.bak"
    except Exception as e:
        return f"Rollback error: {e}"

def restart_bot() -> str:
    try:
        import threading
        def _restart():
            import time
            time.sleep(1.5)
            os.execv(sys.executable, [sys.executable] + sys.argv)
        threading.Thread(target=_restart, daemon=True).start()
        return "Restarting in 1.5 seconds..."
    except Exception as e:
        return f"Restart error: {e}"

def pip_install(package: str) -> str:
    try:
        result = subprocess.run(
            [sys.executable, "-m", "pip", "install", package],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode == 0:
            return f"Installed: {package}"
        return f"pip error: {result.stderr[-500:]}"
    except Exception as e:
        return f"Error: {e}"

def download_file(url: str, save_path: str) -> str:
    try:
        os.makedirs(os.path.dirname(save_path), exist_ok=True) if os.path.dirname(save_path) else None
        resp = requests.get(url, stream=True, timeout=60, headers={"User-Agent": "Mozilla/5.0"})
        resp.raise_for_status()
        with open(save_path, "wb") as f:
            for chunk in resp.iter_content(chunk_size=8192):
                f.write(chunk)
        size_mb = os.path.getsize(save_path) / (1024 * 1024)
        return f"Downloaded: {save_path} ({size_mb:.1f} MB)"
    except Exception as e:
        return f"Download error: {e}"

def teams_launch_app() -> str:
    """Launch Teams desktop app via PowerShell."""
    try:
        r = subprocess.run(
            ["powershell", "-Command", "Start-Process 'shell:AppsFolder\\MSTeams_8wekyb3d8bbwe!MSTeams'"],
            capture_output=True, text=True, timeout=10
        )
        if r.returncode == 0:
            return "Teams launched."
        return f"Error: {r.stderr.strip()}"
    except Exception as e:
        return f"Error launching Teams: {e}"



async def _teams_page():
    """Get browser page, assumed already logged into Teams."""
    page, err = await _get_page()
    return page, err

async def teams_login() -> str:
    creds = _load_creds().get("teams", {})
    email = creds.get("email", "")
    password = creds.get("password", "")
    if not email or not password:
        return (
            "Teams credentials not set. Tell Adam:\n"
            "save my teams email as your@email.com\n"
            "save my teams password as yourpassword"
        )
    page, err = await _get_page()
    if err:
        return err
    try:
        await page.goto("https://teams.microsoft.com", wait_until="domcontentloaded", timeout=30000)
        await page.wait_for_timeout(3000)

        # Already logged in?
        if "teams.microsoft.com/_" in page.url or "teams.live.com" in page.url:
            return "Already logged into Teams."

        # Fill email
        try:
            await page.fill('input[type="email"]', email, timeout=8000)
            await page.click('input[type="submit"]', timeout=5000)
            await page.wait_for_timeout(2000)
        except Exception:
            pass

        # Fill password
        try:
            await page.fill('input[type="password"]', password, timeout=8000)
            await page.click('input[type="submit"]', timeout=5000)
            await page.wait_for_timeout(3000)
        except Exception:
            pass

        # MFA page — always pick first option (Microsoft Authenticator)
        try:
            # Detect MFA screen by looking for the authenticator option
            mfa_option = page.locator("text=Authenticator").first
            if await mfa_option.count() > 0:
                await mfa_option.click(timeout=5000)
            else:
                # Try clicking the first list item on the MFA page
                first_option = page.locator('[data-value], .table, li').first
                if await first_option.count() > 0:
                    await first_option.click(timeout=5000)
        except Exception:
            pass

        # Wait for phone approval (up to 60 seconds)
        try:
            await page.wait_for_url("**/teams.microsoft.com/**", timeout=60000)
        except Exception:
            pass

        # "Stay signed in?" prompt
        try:
            await page.click('input[type="submit"]', timeout=5000)
            await page.wait_for_timeout(2000)
        except Exception:
            pass

        # Wait for Teams to fully load
        await page.wait_for_timeout(5000)
        title = await page.title()
        return f"Teams opened: {title} — approve the Authenticator notification on your phone if prompted."
    except Exception as e:
        return f"Teams login error: {e}"

async def teams_get_chats() -> str:
    page, err = await _teams_page()
    if err:
        return err
    try:
        # Navigate to chat section
        await page.wait_for_timeout(2000)
        chat_items = await page.evaluate("""() => {
            const items = [];
            document.querySelectorAll('[data-tid="chat-list-item"], [role="listitem"]').forEach(el => {
                const name = el.querySelector('[class*="title"], [class*="name"]')?.innerText?.trim();
                const preview = el.querySelector('[class*="preview"], [class*="subtitle"]')?.innerText?.trim();
                if (name) items.push({ name, preview: preview || '' });
            });
            return items.slice(0, 15);
        }""")
        if not chat_items:
            return "No chats found. Make sure Teams is open and logged in (run teams_login first)."
        lines = [f"{i+1}. {c['name']} — {c['preview']}" for i, c in enumerate(chat_items)]
        return "\n".join(lines)
    except Exception as e:
        return f"Error reading chats: {e}"

async def teams_open_chat(name: str) -> str:
    page, err = await _teams_page()
    if err:
        return err
    try:
        # Try clicking on the chat by name
        item = page.get_by_text(name, exact=False).first
        await item.click(timeout=8000)
        await page.wait_for_timeout(2000)
        return f"Opened chat: {name}"
    except Exception as e:
        return f"Could not open chat '{name}': {e}"

async def teams_read_messages() -> str:
    page, err = await _teams_page()
    if err:
        return err
    try:
        await page.wait_for_timeout(1500)
        messages = await page.evaluate("""() => {
            const msgs = [];
            document.querySelectorAll('[data-tid*="message"], [class*="messageBody"], [class*="message-body"]').forEach(el => {
                const text = el.innerText?.trim();
                if (text && text.length > 1) msgs.push(text);
            });
            return msgs.slice(-20);
        }""")
        if not messages:
            # Fallback: grab all visible text in the message thread
            content = await page.inner_text('[id*="message-list"], [class*="messageThread"], [role="list"]')
            return content[:3000] if content else "No messages found in current chat."
        return "\n---\n".join(messages)
    except Exception as e:
        return f"Error reading messages: {e}"

async def teams_send_message(message: str) -> str:
    page, err = await _teams_page()
    if err:
        return err
    try:
        # Find message input box
        box = page.locator('[data-tid="ckeditor"], [contenteditable="true"][role="textbox"], [id*="message-input"]').first
        await box.click(timeout=5000)
        await box.type(message, delay=30)
        await page.wait_for_timeout(500)
        # Send with Enter
        await page.keyboard.press("Enter")
        await page.wait_for_timeout(1000)
        return f"Message sent: {message}"
    except Exception as e:
        return f"Error sending message: {e}"

async def teams_search(query: str) -> str:
    page, err = await _teams_page()
    if err:
        return err
    try:
        search = page.locator('[data-tid="app-bar-search"], input[placeholder*="Search"], [aria-label*="Search"]').first
        await search.click(timeout=5000)
        await search.fill(query)
        await page.wait_for_timeout(2000)
        content = await page.inner_text('[class*="searchResult"], [data-tid*="search"]')
        return content[:2000] if content else "No results found."
    except Exception as e:
        return f"Search error: {e}"


async def teams_list_classes() -> str:
    page, err = await _teams_page()
    if err:
        return err
    try:
        await page.wait_for_timeout(2000)
        classes = await page.evaluate("""() => {
            const items = [];
            // Teams sidebar items
            document.querySelectorAll('[data-tid="team-channel-list-team-title"], [class*="teamName"], [aria-label*="team"]').forEach(el => {
                const t = el.innerText?.trim();
                if (t) items.push(t);
            });
            return [...new Set(items)].slice(0, 20);
        }""")
        if not classes:
            # Fallback: get all sidebar text
            sidebar = await page.inner_text('[class*="teamList"], [class*="sidebar"], nav')
            return sidebar[:2000] if sidebar else "No classes found. Make sure Teams web is open."
        return "\n".join(f"{i+1}. {c}" for i, c in enumerate(classes))
    except Exception as e:
        return f"Error listing classes: {e}"


async def teams_open_class(name: str) -> str:
    page, err = await _teams_page()
    if err:
        return err
    try:
        # Click on the team/class in the sidebar
        item = page.get_by_text(name, exact=False).first
        await item.click(timeout=8000)
        await page.wait_for_timeout(2000)
        return f"Opened class: {name}"
    except Exception as e:
        return f"Could not open class '{name}': {e}"


async def teams_get_assignments() -> str:
    page, err = await _teams_page()
    if err:
        return err
    try:
        # Click "Assignments" tab
        tab = page.get_by_role("tab", name=lambda n: "ssignment" in n).first
        if await tab.count() == 0:
            tab = page.get_by_text("Assignments", exact=False).first
        await tab.click(timeout=8000)
        await page.wait_for_timeout(3000)

        assignments = await page.evaluate("""() => {
            const items = [];
            document.querySelectorAll('[class*="assignment"], [data-tid*="assignment"]').forEach(el => {
                const title = el.querySelector('[class*="title"], h3, h2')?.innerText?.trim();
                const due = el.querySelector('[class*="due"], [class*="date"]')?.innerText?.trim();
                if (title) items.push(title + (due ? ' | Due: ' + due : ''));
            });
            return items.slice(0, 20);
        }""")

        if not assignments:
            content = await page.inner_text('main, [role="main"]')
            return content[:3000] if content else "No assignments found."
        return "\n".join(f"{i+1}. {a}" for i, a in enumerate(assignments))
    except Exception as e:
        return f"Error getting assignments: {e}"


async def teams_open_assignment(name: str) -> str:
    page, err = await _teams_page()
    if err:
        return err
    try:
        item = page.get_by_text(name, exact=False).first
        await item.click(timeout=8000)
        await page.wait_for_timeout(3000)
        return f"Opened assignment: {name}"
    except Exception as e:
        return f"Could not open assignment '{name}': {e}"


async def teams_read_assignment_details() -> str:
    page, err = await _teams_page()
    if err:
        return err
    try:
        await page.wait_for_timeout(2000)
        # Try to grab assignment detail panel content
        content = await page.evaluate("""() => {
            const selectors = [
                '[class*="assignmentDetail"]',
                '[class*="assignment-detail"]',
                '[class*="instructions"]',
                '[role="main"]',
                'main'
            ];
            for (const sel of selectors) {
                const el = document.querySelector(sel);
                if (el && el.innerText?.trim().length > 50) return el.innerText.trim();
            }
            return document.body.innerText.trim();
        }""")
        return content[:5000] if content else "Could not read assignment details."
    except Exception as e:
        return f"Error reading assignment: {e}"


async def teams_download_assignment_file() -> str:
    page, err = await _teams_page()
    if err:
        return err
    try:
        download_dir = os.path.join(os.path.expanduser("~"), "Desktop", "Teams Downloads")
        os.makedirs(download_dir, exist_ok=True)

        # Set up download listener
        async with page.expect_download(timeout=15000) as dl_info:
            # Try clicking any download/file link in the assignment
            file_link = page.locator(
                'a[href*=".docx"], a[href*=".pdf"], a[href*=".pptx"], '
                '[class*="attachment"], [data-tid*="file"], button:has-text("Download")'
            ).first
            await file_link.click(timeout=8000)

        download = await dl_info.value
        save_path = os.path.join(download_dir, download.suggested_filename)
        await download.save_as(save_path)
        return f"Downloaded: {save_path}"
    except Exception as e:
        return f"Download error: {e}"


def send_telegram_file(path: str) -> str:
    path = path.strip().replace('"', '').replace("'", "")
    if os.path.isdir(path):
        import zipfile
        zip_path = path.rstrip("/\\") + "__adam.zip"
        try:
            with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
                for root, dirs, files in os.walk(path):
                    dirs[:] = [d for d in dirs if not d.startswith(".")]
                    for file in files:
                        fp = os.path.join(root, file)
                        arcname = os.path.relpath(fp, os.path.dirname(path))
                        zf.write(fp, arcname)
            size_mb = os.path.getsize(zip_path) / (1024 * 1024)
            if size_mb > 49:
                os.remove(zip_path)
                return f"Folder too large to send ({size_mb:.1f} MB). Telegram limit is 50 MB. Try a specific subfolder or file."
            return f"FILE_TO_SEND:{zip_path}"
        except Exception as e:
            return f"Error zipping folder: {e}"
    elif os.path.isfile(path):
        size_mb = os.path.getsize(path) / (1024 * 1024)
        if size_mb > 49:
            return f"File too large ({size_mb:.1f} MB). Telegram limit is 50 MB."
        return f"FILE_TO_SEND:{path}"
    else:
        return f"Not found: {path}"


# --- Browser (Playwright async) ---

_pw_instance = None
_pw_browser = None
_pw_page = None


async def _get_page():
    global _pw_instance, _pw_browser, _pw_page
    try:
        from playwright.async_api import async_playwright
    except ImportError:
        return None, "Playwright not installed. Run: pip install playwright && python -m playwright install chromium"
    if _pw_page is None or _pw_browser is None or not _pw_browser.is_connected():
        if _pw_instance is None:
            _pw_instance = await async_playwright().start()
        # Use real Chrome with user profile — avoids bot detection, keeps YouTube logged in
        try:
            _pw_browser = await _pw_instance.chromium.launch(
                headless=False,
                channel="chrome",
                args=["--start-maximized", "--disable-blink-features=AutomationControlled"],
            )
        except Exception:
            # Fallback to bundled Chromium if Chrome not found
            _pw_browser = await _pw_instance.chromium.launch(
                headless=False,
                args=["--start-maximized", "--disable-blink-features=AutomationControlled"],
            )
        _pw_page = await _pw_browser.new_page()
        # Hide automation fingerprint
        await _pw_page.add_init_script("Object.defineProperty(navigator, 'webdriver', {get: () => undefined})")
        await _pw_page.set_viewport_size({"width": 1280, "height": 800})
    return _pw_page, None


async def browser_navigate(url: str) -> str:
    page, err = await _get_page()
    if err:
        return err
    try:
        await page.goto(url, wait_until="domcontentloaded", timeout=20000)
        return f"Navigated to: {await page.title()} — {page.url}"
    except Exception as e:
        return f"Error navigating: {e}"


async def browser_click(selector: str) -> str:
    page, err = await _get_page()
    if err:
        return err
    try:
        if selector.startswith("text="):
            await page.get_by_text(selector[5:], exact=False).first.click(timeout=8000)
        else:
            await page.click(selector, timeout=8000)
        return f"Clicked: {selector}"
    except Exception as e:
        return f"Click error: {e}"


async def browser_type(selector: str, text: str, press_enter: bool = False) -> str:
    page, err = await _get_page()
    if err:
        return err
    try:
        await page.fill(selector, text, timeout=8000)
        if press_enter:
            await page.press(selector, "Enter")
        return f"Typed '{text}' into {selector}" + (" + Enter" if press_enter else "")
    except Exception as e:
        return f"Type error: {e}"


async def browser_screenshot() -> str:
    page, err = await _get_page()
    if err:
        return err
    try:
        path = os.path.join(tempfile.gettempdir(), "adam_screenshot.png")
        await page.screenshot(path=path, full_page=False)
        return f"Screenshot saved: {path}"
    except Exception as e:
        return f"Screenshot error: {e}"


async def browser_get_content() -> str:
    page, err = await _get_page()
    if err:
        return err
    try:
        text = await page.inner_text("body")
        return text[:4000]
    except Exception as e:
        return f"Content error: {e}"


async def browser_press_key(key: str) -> str:
    page, err = await _get_page()
    if err:
        return err
    try:
        await page.keyboard.press(key)
        return f"Pressed: {key}"
    except Exception as e:
        return f"Key error: {e}"


async def browser_scroll(direction: str, amount: int = 500) -> str:
    page, err = await _get_page()
    if err:
        return err
    try:
        delta = amount if direction == "down" else -amount
        await page.evaluate(f"window.scrollBy(0, {delta})")
        return f"Scrolled {direction} {amount}px"
    except Exception as e:
        return f"Scroll error: {e}"


async def browser_wait(selector: str = "", ms: int = 1000) -> str:
    page, err = await _get_page()
    if err:
        return err
    try:
        if selector:
            await page.wait_for_selector(selector, timeout=15000)
            return f"Element appeared: {selector}"
        else:
            await page.wait_for_timeout(ms)
            return f"Waited {ms}ms"
    except Exception as e:
        return f"Wait error: {e}"


async def browser_evaluate(script: str) -> str:
    page, err = await _get_page()
    if err:
        return err
    try:
        result = await page.evaluate(script)
        return str(result)[:2000]
    except Exception as e:
        return f"Eval error: {e}"


async def youtube_play(query: str) -> str:
    page, err = await _get_page()
    if err:
        return err
    try:
        import urllib.parse
        url = f"https://www.youtube.com/results?search_query={urllib.parse.quote(query)}"
        await page.goto(url, wait_until="domcontentloaded", timeout=20000)
        await page.wait_for_timeout(2000)

        # Click first non-short video result
        first = page.locator("a#video-title").first
        title = await first.get_attribute("title") or query
        await first.click(timeout=8000)

        # Wait for video player to load
        await page.wait_for_timeout(3000)

        # Skip ad if present (retry a few times as ad button may appear late)
        for _ in range(5):
            try:
                skip = page.locator(".ytp-skip-ad-button, .ytp-ad-skip-button-container button")
                if await skip.count() > 0:
                    await skip.first.click()
                    await page.wait_for_timeout(1000)
                    break
            except Exception:
                pass
            await page.wait_for_timeout(1000)

        # Ensure video is playing (click play if paused)
        try:
            video = page.locator("video")
            paused = await video.evaluate("v => v.paused")
            if paused:
                await video.click()
        except Exception:
            pass

        return f"Playing on YouTube: {title}"
    except Exception as e:
        return f"YouTube play error: {e}"


async def browser_get_interactive_elements() -> str:
    page, err = await _get_page()
    if err:
        return err
    try:
        elements = await page.evaluate("""() => {
            const out = [];
            document.querySelectorAll('input, textarea, [contenteditable="true"]').forEach(el => {
                out.push('INPUT: placeholder="' + (el.placeholder||'') + '" aria-label="' + (el.getAttribute('aria-label')||'') + '" id="' + (el.id||'') + '"');
            });
            document.querySelectorAll('button, [role="button"], a[href]').forEach(el => {
                const txt = (el.innerText || el.getAttribute('aria-label') || el.title || '').trim().slice(0, 60);
                if (txt) out.push('CLICKABLE: "' + txt + '" id="' + (el.id||'') + '"');
            });
            return out.slice(0, 40);
        }""")
        if not elements:
            return "No interactive elements found on page."
        return "\n".join(elements)
    except Exception as e:
        return f"Error: {e}"


async def browser_close() -> str:
    global _pw_instance, _pw_browser, _pw_page
    try:
        if _pw_browser:
            await _pw_browser.close()
        if _pw_instance:
            await _pw_instance.stop()
        _pw_browser = None
        _pw_page = None
        _pw_instance = None
        return "Browser closed."
    except Exception as e:
        return f"Close error: {e}"


def set_volume(level: int) -> str:
    level = max(0, min(100, level))
    try:
        from pycaw.pycaw import AudioUtilities
        vol = AudioUtilities.GetSpeakers().EndpointVolume
        vol.SetMasterVolumeLevelScalar(level / 100.0, None)
        return f"Volume set to {level}%"
    except Exception as e:
        return f"Error setting volume: {e}"


def get_volume() -> str:
    try:
        from pycaw.pycaw import AudioUtilities
        vol = AudioUtilities.GetSpeakers().EndpointVolume
        level = round(vol.GetMasterVolumeLevelScalar() * 100)
        return f"Current volume: {level}%"
    except Exception as e:
        return f"Error reading volume: {e}"


def set_brightness(level: int) -> str:
    level = max(0, min(100, level))
    try:
        import screen_brightness_control as sbc
        sbc.set_brightness(level)
        return f"Brightness set to {level}%"
    except Exception as e:
        return f"Error setting brightness: {e}"


def get_brightness() -> str:
    try:
        import screen_brightness_control as sbc
        b = sbc.get_brightness()
        return f"Current brightness: {b}%"
    except Exception as e:
        return f"Error reading brightness: {e}"


def send_email(to: str, subject: str, body: str, cc: str = "", attachments: list = None) -> str:
    addr, pwd = _gmail_creds()
    if not addr or not pwd:
        return "Gmail credentials not configured."
    try:
        msg = MIMEMultipart()
        msg["From"] = addr
        msg["To"] = to
        msg["Subject"] = subject
        if cc:
            msg["Cc"] = cc
        msg.attach(MIMEText(body, "plain", "utf-8"))

        attached = []
        failed = []
        for path in (attachments or []):
            path = path.strip()
            if not os.path.exists(path):
                failed.append(f"{path} (not found)")
                continue
            try:
                with open(path, "rb") as f:
                    part = MIMEBase("application", "octet-stream")
                    part.set_payload(f.read())
                encoders.encode_base64(part)
                filename = os.path.basename(path)
                part.add_header("Content-Disposition", f'attachment; filename="{filename}"')
                msg.attach(part)
                attached.append(filename)
            except Exception as e:
                failed.append(f"{path} ({e})")

        recipients = [to] + ([cc] if cc else [])
        with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
            server.login(addr, pwd)
            server.sendmail(addr, recipients, msg.as_string())

        summary = f"Email sent to {to} — Subject: {subject}"
        if attached:
            summary += f"\nAttached: {', '.join(attached)}"
        if failed:
            summary += f"\nFailed to attach: {', '.join(failed)}"
        return summary
    except Exception as e:
        return f"Error sending email: {e}"


def media_control(action: str) -> str:
    """Control media playback using Windows media keys."""
    try:
        import win32api, win32con
        VK_MEDIA_PLAY_PAUSE = 0xB3
        VK_MEDIA_NEXT_TRACK = 0xB0
        VK_MEDIA_PREV_TRACK = 0xB1
        VK_MEDIA_STOP       = 0xB2
        key_map = {
            "play_pause": VK_MEDIA_PLAY_PAUSE,
            "next":       VK_MEDIA_NEXT_TRACK,
            "previous":   VK_MEDIA_PREV_TRACK,
            "stop":       VK_MEDIA_STOP,
        }
        if action not in key_map:
            return f"Invalid action. Use: {', '.join(key_map.keys())}"
        vk = key_map[action]
        win32api.keybd_event(vk, 0, 0, 0)
        win32api.keybd_event(vk, 0, win32con.KEYEVENTF_KEYUP, 0)
        return f"Media {action.replace('_',' ')} sent."
    except Exception as e:
        return f"Media control error: {e}"


def create_python_file(path: str, code: str) -> str:
    try:
        os.makedirs(os.path.dirname(path), exist_ok=True) if os.path.dirname(path) else None
        with open(path, "w", encoding="utf-8") as f:
            f.write(code)
        return f"Python file created: {path}"
    except Exception as e:
        return f"Error creating Python file: {e}"


def create_jupyter_notebook(path: str, cells: list) -> str:
    try:
        import json as _json
        nb_cells = []
        for cell in cells:
            ctype = cell.get("type", "code")
            source = cell.get("content", "")
            if ctype == "markdown":
                nb_cells.append({
                    "cell_type": "markdown",
                    "metadata": {},
                    "source": source.splitlines(keepends=True),
                })
            else:
                nb_cells.append({
                    "cell_type": "code",
                    "execution_count": None,
                    "metadata": {},
                    "outputs": [],
                    "source": source.splitlines(keepends=True),
                })
        notebook = {
            "nbformat": 4,
            "nbformat_minor": 5,
            "metadata": {
                "kernelspec": {"display_name": "Python 3", "language": "python", "name": "python3"},
                "language_info": {"name": "python", "version": "3.10.0"},
            },
            "cells": nb_cells,
        }
        os.makedirs(os.path.dirname(path), exist_ok=True) if os.path.dirname(path) else None
        with open(path, "w", encoding="utf-8") as f:
            _json.dump(notebook, f, indent=1)
        return f"Jupyter notebook created: {path}"
    except Exception as e:
        return f"Error creating notebook: {e}"


def create_powerpoint(path: str, slides: list, title: str = "") -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        # Title slide
        if title:
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title
            if slide.placeholders[1:]:
                slide.placeholders[1].text = ""
        # Content slides
        for s in slides:
            layout = prs.slide_layouts[1]  # title + content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = s.get("title", "")
            body = slide.placeholders[1]
            body.text = s.get("content", "")
        os.makedirs(os.path.dirname(path), exist_ok=True) if os.path.dirname(path) else None
        prs.save(path)
        return f"PowerPoint created: {path}"
    except Exception as e:
        return f"Error creating PowerPoint: {e}"


def create_text_file(path: str, content: str) -> str:
    try:
        os.makedirs(os.path.dirname(path), exist_ok=True) if os.path.dirname(path) else None
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        return f"File created: {path}"
    except Exception as e:
        return f"Error creating file: {e}"


def create_word_document(path: str, content: str, title: str = "") -> str:
    try:
        doc = Document()
        if title:
            doc.add_heading(title, 0)
        for line in content.split("\n"):
            doc.add_paragraph(line)
        doc.save(path)
        return f"Word document created: {path}"
    except Exception as e:
        return f"Error creating Word doc: {e}"


def create_pdf(path: str, content: str, title: str = "") -> str:
    try:
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Helvetica", size=12)
        if title:
            pdf.set_font("Helvetica", "B", 16)
            pdf.cell(0, 10, title, ln=True)
            pdf.set_font("Helvetica", size=12)
            pdf.ln(4)
        for line in content.split("\n"):
            pdf.multi_cell(0, 8, line)
        pdf.output(path)
        return f"PDF created: {path}"
    except Exception as e:
        return f"Error creating PDF: {e}"


def create_latex_document(path: str, latex_source: str) -> str:
    try:
        tex_path = path if path.endswith(".tex") else path + ".tex"
        pdf_path = tex_path.replace(".tex", ".pdf")
        tex_dir = os.path.dirname(tex_path) or "."
        with open(tex_path, "w", encoding="utf-8") as f:
            f.write(latex_source)
        # Try xelatex first (better unicode), fall back to pdflatex
        for engine in ("xelatex", "pdflatex"):
            result = subprocess.run(
                [engine, "-interaction=nonstopmode", "-output-directory", tex_dir, tex_path],
                capture_output=True, text=True, timeout=60,
            )
            if result.returncode == 0 and os.path.exists(pdf_path):
                # Run twice for TOC/refs
                subprocess.run(
                    [engine, "-interaction=nonstopmode", "-output-directory", tex_dir, tex_path],
                    capture_output=True, timeout=60,
                )
                return f"LaTeX compiled successfully: {pdf_path}"
            if "not found" not in result.stderr.lower() and result.returncode != 127:
                # Engine found but compile error
                return f"LaTeX compile error ({engine}):\n{result.stdout[-1500:]}"
        # No LaTeX engine installed — return .tex path
        return f"No LaTeX engine found. Source saved: {tex_path}. Install MiKTeX or TeX Live to compile."
    except Exception as e:
        return f"Error creating LaTeX document: {e}"


def create_formal_word_document(path: str, title: str, sections: list,
                                 author: str = "", date: str = "", abstract: str = "") -> str:
    try:
        from docx import Document
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        NAVY    = RGBColor(0x1B, 0x2A, 0x4A)
        STEEL   = RGBColor(0x2E, 0x5F, 0x8A)
        ACCENT  = RGBColor(0x2E, 0x75, 0xB6)
        GRAY    = RGBColor(0x59, 0x59, 0x59)
        WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

        doc = Document()

        # --- Global base style ---
        normal = doc.styles["Normal"]
        normal.font.name = "Calibri"
        normal.font.size = Pt(11)
        normal.paragraph_format.space_after = Pt(6)

        # --- Heading styles ---
        h1s = doc.styles["Heading 1"]
        h1s.font.name = "Calibri Light"; h1s.font.size = Pt(16)
        h1s.font.bold = True; h1s.font.color.rgb = NAVY
        h1s.paragraph_format.space_before = Pt(18); h1s.paragraph_format.space_after = Pt(6)
        h1s.paragraph_format.keep_with_next = True

        h2s = doc.styles["Heading 2"]
        h2s.font.name = "Calibri Light"; h2s.font.size = Pt(13)
        h2s.font.bold = True; h2s.font.color.rgb = STEEL
        h2s.paragraph_format.space_before = Pt(12); h2s.paragraph_format.space_after = Pt(4)
        h2s.paragraph_format.keep_with_next = True

        h3s = doc.styles["Heading 3"]
        h3s.font.name = "Calibri"; h3s.font.size = Pt(11)
        h3s.font.bold = True; h3s.font.italic = True; h3s.font.color.rgb = GRAY
        h3s.paragraph_format.space_before = Pt(8); h3s.paragraph_format.space_after = Pt(3)

        # --- Page margins ---
        for sec in doc.sections:
            sec.top_margin    = Inches(1)
            sec.bottom_margin = Inches(1)
            sec.left_margin   = Inches(1.2)
            sec.right_margin  = Inches(1.2)
            sec.different_first_page_header_footer = True

        # --- XML helpers ---
        def _set_shading(para, fill_hex):
            pPr = para._p.get_or_add_pPr()
            shd = OxmlElement("w:shd")
            shd.set(qn("w:val"), "clear")
            shd.set(qn("w:color"), "auto")
            shd.set(qn("w:fill"), fill_hex)
            pPr.append(shd)

        def _add_bottom_border(para, color_hex="1B2A4A", sz="8"):
            pPr = para._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            bot = OxmlElement("w:bottom")
            bot.set(qn("w:val"), "single"); bot.set(qn("w:sz"), sz)
            bot.set(qn("w:space"), "1"); bot.set(qn("w:color"), color_hex)
            pBdr.append(bot); pPr.append(pBdr)

        def _add_top_border(para, color_hex="BBBBBB", sz="4"):
            pPr = para._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            top = OxmlElement("w:top")
            top.set(qn("w:val"), "single"); top.set(qn("w:sz"), sz)
            top.set(qn("w:space"), "1"); top.set(qn("w:color"), color_hex)
            pBdr.append(top); pPr.append(pBdr)

        def _add_left_border(para, color_hex="2E75B6", sz="24"):
            pPr = para._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            left = OxmlElement("w:left")
            left.set(qn("w:val"), "single"); left.set(qn("w:sz"), sz)
            left.set(qn("w:space"), "4"); left.set(qn("w:color"), color_hex)
            pBdr.append(left); pPr.append(pBdr)

        def _page_number_field(run):
            for tag, text in [("w:fldChar", None), ("w:instrText", "PAGE"), ("w:fldChar", None)]:
                el = OxmlElement(tag)
                if tag == "w:fldChar":
                    el.set(qn("w:fldCharType"), "begin" if text is None and not run._r.findall(f"{{{qn('w:fldChar').split('}')[0][1:]}}}fldChar") else "end")
                else:
                    el.text = text
                run._r.append(el)
            # Simpler approach:
            fldChar1 = OxmlElement("w:fldChar"); fldChar1.set(qn("w:fldCharType"), "begin")
            instrText = OxmlElement("w:instrText"); instrText.text = "PAGE"
            fldChar2 = OxmlElement("w:fldChar"); fldChar2.set(qn("w:fldCharType"), "end")
            run._r.clear(); run._r.append(fldChar1); run._r.append(instrText); run._r.append(fldChar2)

        # ===== COVER PAGE =====
        # Navy banner block
        banner = doc.add_paragraph()
        _set_shading(banner, "1B2A4A")
        banner.paragraph_format.space_before = Pt(0)
        banner.paragraph_format.space_after  = Pt(0)
        pad_run = banner.add_run("\n\n")
        pad_run.font.size = Pt(20)
        pad_run.font.color.rgb = WHITE
        title_run = banner.add_run(f"{title}\n")
        title_run.font.name = "Calibri Light"; title_run.font.size = Pt(28)
        title_run.bold = True; title_run.font.color.rgb = WHITE
        pad2 = banner.add_run("\n")
        pad2.font.size = Pt(14); pad2.font.color.rgb = WHITE

        # Thin accent stripe
        stripe = doc.add_paragraph()
        _set_shading(stripe, "2E75B6")
        stripe.paragraph_format.space_before = Pt(0)
        stripe.paragraph_format.space_after  = Pt(0)
        stripe.add_run(" ").font.size = Pt(3)

        # Spacer
        for _ in range(5):
            sp = doc.add_paragraph(); sp.paragraph_format.space_after = Pt(0)

        if author:
            ap = doc.add_paragraph()
            ar = ap.add_run(author)
            ar.font.name = "Calibri"; ar.font.size = Pt(14)
            ar.bold = True; ar.font.color.rgb = NAVY
            ap.paragraph_format.space_after = Pt(3)

        if date:
            dp = doc.add_paragraph()
            dr = dp.add_run(date)
            dr.font.name = "Calibri"; dr.font.size = Pt(11)
            dr.font.color.rgb = GRAY
            dp.paragraph_format.space_after = Pt(0)

        doc.add_page_break()

        # ===== ABSTRACT =====
        if abstract:
            h = doc.add_heading("Abstract", level=1)
            _add_bottom_border(h)
            abs_p = doc.add_paragraph()
            abs_p.paragraph_format.left_indent  = Inches(0.3)
            abs_p.paragraph_format.right_indent = Inches(0.3)
            abs_p.paragraph_format.space_after  = Pt(16)
            _add_left_border(abs_p, "2E75B6", "20")
            ar = abs_p.add_run(abstract)
            ar.font.name = "Calibri"; ar.font.size = Pt(10.5)
            ar.italic = True; ar.font.color.rgb = GRAY
            doc.add_page_break()

        # ===== BODY SECTIONS =====
        for sec in sections:
            heading_text = sec.get("heading", "")
            content      = sec.get("content", "")
            level        = min(int(sec.get("level", 1)), 3)

            if heading_text:
                h = doc.add_heading(heading_text, level=level)
                if level == 1:
                    _add_bottom_border(h, "1B2A4A", "6")

            for block in (content or "").split("\n\n"):
                block = block.strip()
                if not block:
                    continue
                lines = [l for l in block.split("\n") if l.strip()]
                is_bullets = all(l.strip()[0] in "-•*" for l in lines)
                if is_bullets:
                    for line in lines:
                        line = line.lstrip("-•* ").strip()
                        if not line:
                            continue
                        bp = doc.add_paragraph(style="List Bullet")
                        br = bp.add_run(line)
                        br.font.name = "Calibri"; br.font.size = Pt(11)
                        bp.paragraph_format.space_after = Pt(2)
                else:
                    p = doc.add_paragraph()
                    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                    p.paragraph_format.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
                    p.paragraph_format.line_spacing = 1.15
                    p.paragraph_format.space_after  = Pt(8)
                    pr = p.add_run(block)
                    pr.font.name = "Calibri"; pr.font.size = Pt(11)

        # ===== HEADER (skips cover via different_first_page) =====
        for sec in doc.sections:
            hdr = sec.header
            hp  = hdr.paragraphs[0] if hdr.paragraphs else hdr.add_paragraph()
            hp.clear(); hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            hr = hp.add_run(title[:70])
            hr.font.name = "Calibri"; hr.font.size = Pt(9)
            hr.italic = True; hr.font.color.rgb = GRAY
            _add_bottom_border(hp, "BBBBBB", "4")

        # ===== FOOTER =====
        for sec in doc.sections:
            ftr = sec.footer
            fp  = ftr.paragraphs[0] if ftr.paragraphs else ftr.add_paragraph()
            fp.clear(); fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            _add_top_border(fp, "BBBBBB", "4")
            if author:
                lr = fp.add_run(f"{author}  |  ")
                lr.font.name = "Calibri"; lr.font.size = Pt(9)
                lr.font.color.rgb = GRAY
            pnr = fp.add_run()
            _page_number_field(pnr)
            pnr.font.name = "Calibri"; pnr.font.size = Pt(9)
            pnr.font.color.rgb = GRAY

        doc.save(path)
        return f"Formal Word document created: {path}"
    except Exception as e:
        return f"Error creating formal Word document: {e}"


def create_spreadsheet(path: str, headers: list, rows) -> str:
    try:
        import json
        if isinstance(rows, str):
            rows = json.loads(rows)
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.append(headers)
        for row in rows:
            ws.append(row)
        wb.save(path)
        return f"Spreadsheet created: {path}"
    except Exception as e:
        return f"Error creating spreadsheet: {e}"


# ── ISSE certificate system tools ────────────────────────────────────────────

def isse_get_ciclos() -> str:
    try:
        resp = requests.get(f"{ISSE_BASE_URL}/api/ciclos", timeout=10)
        resp.raise_for_status()
        ciclos = resp.json().get("ciclos", [])
        if not ciclos:
            return "No hay ciclos cargados en la base de datos ISSE."
        return "Ciclos disponibles en ISSE:\n" + "\n".join(f"  - {c}" for c in ciclos)
    except Exception as e:
        return f"Error consultando ISSE: {e}"


def isse_search_professor(name: str) -> str:
    try:
        resp = requests.get(f"{ISSE_BASE_URL}/api/professors", params={"q": name}, timeout=10)
        resp.raise_for_status()
        profs = resp.json().get("professors", [])
        if not profs:
            return f"No se encontró ningún profesor con el nombre '{name}' en ISSE."
        lines = [f"  - {p}" for p in profs[:15]]
        return f"Profesores encontrados ({len(profs)}):\n" + "\n".join(lines)
    except Exception as e:
        return f"Error buscando profesor en ISSE: {e}"


def isse_get_certificates(professor_id: str, ciclos: list = None) -> str:
    try:
        params = {"profesor": professor_id}
        if ciclos:
            params["ciclos"] = ",".join(ciclos)
        resp = requests.get(f"{ISSE_BASE_URL}/api/certificates", params=params, timeout=15)
        resp.raise_for_status()
        data = resp.json().get("results", [])
        if not data:
            if ciclos:
                ciclos_str = ", ".join(ciclos)
                return f"El profesor '{professor_id}' no tiene registros para los períodos: {ciclos_str}."
            return f"No se encontraron certificados para '{professor_id}'."
        # Group by period so all ciclos appear regardless of total row count
        from collections import defaultdict
        by_ciclo = defaultdict(list)
        for row in data:
            by_ciclo[row.get('ciclo_lectivo', '')].append(row)

        lines = [f"Certificados de '{professor_id}' — {len(data)} registros en {len(by_ciclo)} período(s):\n"]
        for ciclo in sorted(by_ciclo.keys()):
            rows_c = by_ciclo[ciclo]
            total_hrs = sum(int(r.get('horas_semestre', 0) or 0) for r in rows_c)
            lines.append(f"{ciclo} ({len(rows_c)} cursos, {total_hrs} hrs totales):")
            for r in rows_c:
                curso = r.get('nombre_curso') or r.get('materia', '')
                hrs = int(r.get('horas_semestre', 0) or 0)
                dpto = r.get('departamento', '')
                lines.append(f"  - {curso} — {hrs} hrs — {dpto}")
        return "\n".join(lines)
    except Exception as e:
        return f"Error obteniendo certificados ISSE: {e}"


def isse_export_certificate(professor_name: str, ciclos: list = None, fmt: str = "excel") -> str:
    try:
        params = {"profesor": professor_name, "format": fmt}
        if ciclos:
            params["ciclos"] = ",".join(ciclos)
        resp = requests.get(f"{ISSE_BASE_URL}/api/export", params=params, timeout=30)
        resp.raise_for_status()
        suffix = ".xlsx" if fmt == "excel" else ".csv"
        safe_name = professor_name.replace(" ", "_").replace("/", "-")[:40]
        out_path = os.path.join(r"C:\Users\dani1\Downloads", f"certificado_{safe_name}{suffix}")
        with open(out_path, "wb") as f:
            f.write(resp.content)
        return out_path
    except Exception as e:
        return f"Error exportando certificado ISSE: {e}"


async def execute_tool(name: str, args: dict) -> str:
    if name == "run_command":
        return run_command(args["command"])
    elif name == "read_file":
        return read_file(args["path"], args.get("max_chars", 60000))
    elif name == "write_file":
        return write_file(args["path"], args["content"])
    elif name == "list_directory":
        return list_directory(args["path"])
    elif name == "web_search":
        return web_search(args["query"], args.get("max_results", 5))
    elif name == "fetch_webpage":
        return fetch_webpage(args["url"])
    elif name == "get_emails":
        return get_emails(args.get("count", 5), args.get("folder", "INBOX"))
    elif name == "search_emails":
        return search_emails(args["query"], args.get("count", 5))
    elif name == "save_to_memory":
        return save_to_memory(args["content"])
    elif name == "read_memory":
        return read_memory()
    elif name == "obsidian_read_note":
        return obsidian_read_note(args["note_name"])
    elif name == "obsidian_write_note":
        return obsidian_write_note(args["note_name"], args["content"])
    elif name == "obsidian_append_note":
        return obsidian_append_note(args["note_name"], args["content"])
    elif name == "obsidian_list_notes":
        return obsidian_list_notes(args.get("subfolder", ""))
    elif name == "obsidian_search_notes":
        return obsidian_search_notes(args["query"])
    elif name == "obsidian_create_folder":
        return obsidian_create_folder(args["folder_path"])
    elif name == "obsidian_move_note":
        return obsidian_move_note(args["source"], args["destination"])
    elif name == "obsidian_delete_note":
        return obsidian_delete_note(args["note_path"])
    elif name == "create_python_file":
        return create_python_file(args["path"], args["code"])
    elif name == "create_jupyter_notebook":
        return create_jupyter_notebook(args["path"], args["cells"])
    elif name == "create_powerpoint":
        return create_powerpoint(args["path"], args["slides"], args.get("title", ""))
    elif name == "create_text_file":
        return create_text_file(args["path"], args["content"])
    elif name == "create_word_document":
        return create_word_document(args["path"], args["content"], args.get("title", ""))
    elif name == "create_pdf":
        return create_pdf(args["path"], args["content"], args.get("title", ""))
    elif name == "create_latex_document":
        return create_latex_document(args["path"], args["latex_source"])
    elif name == "create_formal_word_document":
        return create_formal_word_document(
            args["path"], args["title"], args["sections"],
            args.get("author", ""), args.get("date", ""), args.get("abstract", ""),
        )
    elif name == "create_spreadsheet":
        return create_spreadsheet(args["path"], args["headers"], args["rows"])
    elif name == "send_email":
        return send_email(args["to"], args["subject"], args["body"], args.get("cc", ""), args.get("attachments", []))
    elif name == "media_control":
        return media_control(args["action"])
    elif name == "set_volume":
        return set_volume(args["level"])
    elif name == "get_volume":
        return get_volume()
    elif name == "set_brightness":
        return set_brightness(args["level"])
    elif name == "get_brightness":
        return get_brightness()
    elif name == "save_credential":
        return save_credential(args["service"], args["key"], args["value"])
    elif name == "list_credentials":
        return list_credentials()
    elif name == "desktop_screenshot":
        return desktop_screenshot(args.get("window_title", ""))
    elif name == "app_focus_window":
        return app_focus_window(args["title"])
    elif name == "app_click":
        return app_click(args["x"], args["y"], args.get("button", "left"), args.get("double", False))
    elif name == "app_type":
        return app_type(args["text"])
    elif name == "app_hotkey":
        return app_hotkey(args["keys"])
    elif name == "app_scroll":
        return app_scroll(args["x"], args["y"], args["direction"], args.get("clicks", 5))
    elif name == "app_click_text":
        return app_click_text(args["text"], args["window_title"])
    elif name == "read_bot_code":
        return read_bot_code(args["file"])
    elif name == "edit_bot_code":
        return edit_bot_code(args["file"], args["old_string"], args["new_string"])
    elif name == "restart_bot":
        return restart_bot()
    elif name == "rollback_bot_code":
        return rollback_bot_code(args["file"])
    elif name == "read_improvement_log":
        return read_improvement_log(args.get("limit", 20))
    elif name == "write_improvement_log":
        _log_improvement("(sonnet)", args["status"], args["detail"], args.get("snippet", ""))
        return "Improvement attempt logged."
    elif name == "save_pending_task":
        import json as _json
        _pf = os.path.join(os.path.dirname(os.path.abspath(__file__)), "pending_task.json")
        with open(_pf, "w") as _f:
            _json.dump({
                "chat_id": args["chat_id"],
                "task": args["task"],
                "attempts": 0,
                "history": args.get("history", [])  # preserve conversation history
            }, _f)
        return f"Pending task saved with context. Will re-attempt after restart."
    elif name == "pip_install":
        return pip_install(args["package"])
    elif name == "download_file":
        return download_file(args["url"], args["save_path"])
    elif name == "teams_launch_app":
        return teams_launch_app()
    elif name == "send_telegram_file":
        return send_telegram_file(args["path"])
    elif name == "youtube_play":
        return await youtube_play(args["query"])
    elif name == "browser_get_interactive_elements":
        return await browser_get_interactive_elements()
    elif name == "browser_navigate":
        return await browser_navigate(args["url"])
    elif name == "browser_click":
        return await browser_click(args["selector"])
    elif name == "browser_type":
        return await browser_type(args["selector"], args["text"], args.get("press_enter", False))
    elif name == "browser_screenshot":
        return await browser_screenshot()
    elif name == "browser_get_content":
        return await browser_get_content()
    elif name == "browser_press_key":
        return await browser_press_key(args["key"])
    elif name == "browser_scroll":
        return await browser_scroll(args["direction"], args.get("amount", 500))
    elif name == "browser_wait":
        return await browser_wait(args.get("selector", ""), args.get("ms", 1000))
    elif name == "browser_evaluate":
        return await browser_evaluate(args["script"])
    elif name == "browser_close":
        return await browser_close()
    elif name == "isse_get_ciclos":
        return isse_get_ciclos()
    elif name == "isse_search_professor":
        return isse_search_professor(args["name"])
    elif name == "isse_get_certificates":
        return isse_get_certificates(args["professor_id"], ciclos=args.get("ciclos"))
    elif name == "isse_export_certificate":
        return isse_export_certificate(
            args["professor_name"],
            ciclos=args.get("ciclos"),
            fmt=args.get("fmt", "excel"),
        )
    else:
        return f"Unknown tool: {name}"
